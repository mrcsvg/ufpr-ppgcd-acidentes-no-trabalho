"""Auditoria geometrica do deck, no lugar da inspecao visual.

O LibreOffice nao funciona neste ambiente, entao nao da para renderizar os slides
e olhar. Esta auditoria confere por calculo o que a vista pegaria: elemento fora
do slide, margem insuficiente, sobreposicao entre formas e texto que nao cabe na
caixa.

A estimativa de largura do texto usa a razao largura/altura media das fontes
usadas (Calibri e Cambria); e aproximada de proposito, com folga, para acusar o
que estoura de forma clara sem reclamar de cada caso limite.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).resolve().parent / "apresentacao-cat.pptx"

MARGEM_MINIMA = 0.45      # polegadas
# Caixas empilhadas encostam de proposito (valor em cima do rotulo): so acima
# desta folga a sobreposicao e real.
FOLGA_ENTRE_FORMAS = 0.02
# Largura media de caractere como fracao do corpo da fonte.
LARGURA_CARACTERE = {"Calibri": 0.48, "Cambria": 0.50}
ALTURA_LINHA = 1.25       # entrelinha como fracao do corpo


def pol(valor) -> float:
    return Emu(valor).inches if valor is not None else 0.0


def _texto_de(forma) -> tuple[str, float, str]:
    """Devolve (texto, corpo em pt, fonte) da forma, ou ("", 0, "")."""
    if not forma.has_text_frame:
        return "", 0.0, ""
    partes, tamanho, fonte = [], 0.0, "Calibri"
    for paragrafo in forma.text_frame.paragraphs:
        for run in paragrafo.runs:
            partes.append(run.text)
            if run.font.size:
                tamanho = max(tamanho, run.font.size.pt)
            if run.font.name:
                fonte = run.font.name
    return "".join(partes), tamanho, fonte


def linhas_estimadas(texto: str, corpo: float, fonte: str, largura_pol: float) -> int:
    if not texto or not corpo or largura_pol <= 0:
        return 0
    largura_char = LARGURA_CARACTERE.get(fonte, 0.48) * corpo / 72
    por_linha = max(1, int(largura_pol / largura_char))
    linhas = 1
    atual = 0
    for palavra in texto.split():
        acrescimo = len(palavra) + (1 if atual else 0)
        if atual + acrescimo > por_linha:
            linhas += 1
            atual = len(palavra)
        else:
            atual += acrescimo
    return linhas


def _um_contem_o_outro(a: dict, b: dict) -> bool:
    """Indica se uma das caixas envolve a outra por completo."""
    def contem(fora, dentro):
        return (fora["x"] <= dentro["x"] + 0.01
                and fora["y"] <= dentro["y"] + 0.01
                and fora["x"] + fora["w"] >= dentro["x"] + dentro["w"] - 0.01
                and fora["y"] + fora["h"] >= dentro["y"] + dentro["h"] - 0.01)

    return contem(a, b) or contem(b, a)


def auditar() -> int:
    pres = Presentation(DECK)
    larg_slide, alt_slide = pol(pres.slide_width), pol(pres.slide_height)
    total = len(pres.slides._sldIdLst)
    print(f'slide: {larg_slide:.2f}" x {alt_slide:.2f}"  ·  {total} slides\n')

    problemas = 0
    for numero, slide in enumerate(pres.slides, 1):
        caixas = []
        for forma in slide.shapes:
            x, y = pol(forma.left), pol(forma.top)
            w, h = pol(forma.width), pol(forma.height)
            texto, corpo, fonte = _texto_de(forma)
            caixas.append({"forma": forma, "x": x, "y": y, "w": w, "h": h,
                           "texto": texto, "corpo": corpo, "fonte": fonte})

            # 1. dentro do slide, com margem
            if x < MARGEM_MINIMA or y < 0 or x + w > larg_slide - MARGEM_MINIMA + 0.01 \
                    or y + h > alt_slide - 0.15:
                print(f"  [slide {numero}] FORA/MARGEM: {texto[:40]!r} "
                      f"em ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} → direita={x+w:.2f} "
                      f"base={y+h:.2f}")
                problemas += 1

            # 2. o texto cabe na caixa?
            if texto and corpo:
                linhas = linhas_estimadas(texto, corpo, fonte, w)
                altura_necessaria = linhas * corpo * ALTURA_LINHA / 72
                if altura_necessaria > h + 0.02:
                    print(f"  [slide {numero}] TEXTO NAO CABE: {texto[:45]!r} "
                          f"({linhas} linhas ≈ {altura_necessaria:.2f}\" em h={h:.2f}\")")
                    problemas += 1

        # 3. sobreposicao entre elementos que carregam conteudo
        for i, a in enumerate(caixas):
            for b in caixas[i + 1:]:
                if a["w"] * a["h"] == 0 or b["w"] * b["h"] == 0:
                    continue
                # Selos (circulo + numero) sao sobrepostos de proposito.
                if min(a["w"], b["w"]) < 0.4 and min(a["h"], b["h"]) < 0.4:
                    continue
                # Card de fundo com conteudo por cima tambem e proposital: um
                # deles nao tem texto e contem o outro por inteiro.
                if _um_contem_o_outro(a, b) and not (a["texto"] and b["texto"]):
                    continue
                sobre_x = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
                sobre_y = min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"])
                if sobre_x > FOLGA_ENTRE_FORMAS and sobre_y > FOLGA_ENTRE_FORMAS:
                    print(f"  [slide {numero}] SOBREPOSICAO de "
                          f"{sobre_x:.2f}\"x{sobre_y:.2f}\": "
                          f"{a['texto'][:28]!r} × {b['texto'][:28]!r}")
                    problemas += 1

    print(f"\nproblemas encontrados: {problemas}")
    return 1 if problemas else 0


if __name__ == "__main__":
    sys.exit(auditar())
